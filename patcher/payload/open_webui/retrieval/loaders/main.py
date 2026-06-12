import asyncio
import io
import json
import logging
import os
import shutil
import sys

import ftfy
import requests
from azure.identity import DefaultAzureCredential
from langchain_community.document_loaders import (
    AzureAIDocumentIntelligenceLoader,
    BSHTMLLoader,
    CSVLoader,
    Docx2txtLoader,
    OutlookMessageLoader,
    PyPDFLoader,
    TextLoader,
    YoutubeLoader,
)
from langchain_core.documents import Document
from open_webui.env import AIOHTTP_CLIENT_SESSION_SSL, GLOBAL_LOG_LEVEL, REQUESTS_VERIFY
from open_webui.retrieval.loaders.datalab_marker import DatalabMarkerLoader
from open_webui.retrieval.loaders.external_document import ExternalDocumentLoader
from open_webui.retrieval.loaders.mineru import MinerULoader
from open_webui.retrieval.loaders.mistral import MistralLoader
from open_webui.retrieval.loaders.paddleocr_vl import PaddleOCRVLLoader

logging.basicConfig(stream=sys.stdout, level=GLOBAL_LOG_LEVEL)
log = logging.getLogger(__name__)

known_source_ext = [
    'go',
    'py',
    'java',
    'sh',
    'bat',
    'ps1',
    'cmd',
    'js',
    'ts',
    'css',
    'cpp',
    'hpp',
    'h',
    'c',
    'cs',
    'sql',
    'log',
    'ini',
    'pl',
    'pm',
    'r',
    'dart',
    'dockerfile',
    'env',
    'php',
    'hs',
    'hsc',
    'lua',
    'nginxconf',
    'conf',
    'm',
    'mm',
    'plsql',
    'perl',
    'rb',
    'rs',
    'db2',
    'scala',
    'bash',
    'swift',
    'vue',
    'svelte',
    'ex',
    'exs',
    'erl',
    'tsx',
    'jsx',
    'hs',
    'lhs',
    'json',
    'yaml',
    'yml',
    'toml',
]

image_file_ext = ['png', 'jpg', 'jpeg', 'webp', 'bmp', 'gif', 'tif', 'tiff']


class ExcelLoader:
    """Fallback Excel loader using pandas when unstructured is not installed."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> list[Document]:
        import pandas as pd

        text_parts = []
        xls = pd.ExcelFile(self.file_path)
        for sheet_name in xls.sheet_names:
            df = pd.read_excel(xls, sheet_name=sheet_name)
            text_parts.append(f'Sheet: {sheet_name}\n{df.to_string(index=False)}')
        return [
            Document(
                page_content='\n\n'.join(text_parts),
                metadata={'source': self.file_path},
            )
        ]


class PptxLoader:
    """Fallback PowerPoint loader using python-pptx when unstructured is not installed."""

    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> list[Document]:
        from pptx import Presentation

        prs = Presentation(self.file_path)
        text_parts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_texts.append(shape.text_frame.text)
            if slide_texts:
                text_parts.append(f'Slide {i}:\n' + '\n'.join(slide_texts))
        return [
            Document(
                page_content='\n\n'.join(text_parts),
                metadata={'source': self.file_path},
            )
        ]


class TikaLoader:
    def __init__(self, url, file_path, mime_type=None, extract_images=None):
        self.url = url
        self.file_path = file_path
        self.mime_type = mime_type

        self.extract_images = extract_images

    def load(self) -> list[Document]:
        with open(self.file_path, 'rb') as f:
            data = f.read()

        if self.mime_type is not None:
            headers = {'Content-Type': self.mime_type}
        else:
            headers = {}

        if self.extract_images == True:
            headers['X-Tika-PDFextractInlineImages'] = 'true'

        endpoint = self.url
        if not endpoint.endswith('/'):
            endpoint += '/'
        endpoint += 'tika/text'

        r = requests.put(endpoint, data=data, headers=headers, verify=REQUESTS_VERIFY)

        if r.ok:
            raw_metadata = r.json()
            text = raw_metadata.get('X-TIKA:content', '<No text content found>').strip()

            if 'Content-Type' in raw_metadata:
                headers['Content-Type'] = raw_metadata['Content-Type']

            log.debug('Tika extracted text: %s', text)

            return [Document(page_content=text, metadata=headers)]
        else:
            raise Exception(f'Error calling Tika: {r.reason}')


class DoclingLoader:
    def __init__(self, url, api_key=None, file_path=None, mime_type=None, params=None):
        self.url = url.rstrip('/')
        self.api_key = api_key
        self.file_path = file_path
        self.mime_type = mime_type

        self.params = params or {}

    def load(self) -> list[Document]:
        with open(self.file_path, 'rb') as f:
            headers = {}
            if self.api_key:
                headers['X-Api-Key'] = f'{self.api_key}'

            r = requests.post(
                f'{self.url}/v1/convert/file',
                files={
                    'files': (
                        self.file_path,
                        f,
                        self.mime_type or 'application/octet-stream',
                    )
                },
                data={
                    'image_export_mode': 'placeholder',
                    **self.params,
                },
                headers=headers,
                verify=AIOHTTP_CLIENT_SESSION_SSL,
            )
        if r.ok:
            result = r.json()
            document_data = result.get('document', {})
            text = document_data.get('md_content', '<No text content found>')

            metadata = {'Content-Type': self.mime_type} if self.mime_type else {}

            log.debug('Docling extracted text: %s', text)
            return [Document(page_content=text, metadata=metadata)]
        else:
            error_msg = f'Error calling Docling API: {r.reason}'
            if r.text:
                try:
                    error_data = r.json()
                    if 'detail' in error_data:
                        error_msg += f' - {error_data["detail"]}'
                except Exception:
                    error_msg += f' - {r.text}'
            raise Exception(f'Error calling Docling: {error_msg}')


def _resolve_tesseract_cmd() -> str | None:
    env_cmd = os.environ.get('TESSERACT_CMD', '').strip()
    if env_cmd:
        return env_cmd

    path_cmd = shutil.which('tesseract')
    if path_cmd:
        return path_cmd

    common_windows_paths = [
        r'C:\Program Files\Tesseract-OCR\tesseract.exe',
        r'C:\Program Files (x86)\Tesseract-OCR\tesseract.exe',
    ]
    for candidate in common_windows_paths:
        if os.path.exists(candidate):
            return candidate

    return None


def _configure_tesseract() -> bool:
    try:
        import pytesseract
    except ImportError:
        if not getattr(_configure_tesseract, '_missing_dependency_logged', False):
            log.warning('pytesseract is not installed. OCR fallback is unavailable.')
            _configure_tesseract._missing_dependency_logged = True
        return False

    if getattr(_configure_tesseract, '_configured', False):
        return True

    tesseract_cmd = _resolve_tesseract_cmd()
    if not tesseract_cmd:
        if not getattr(_configure_tesseract, '_missing_binary_logged', False):
            log.warning(
                'Tesseract binary was not found. Set TESSERACT_CMD or add tesseract to PATH.'
            )
            _configure_tesseract._missing_binary_logged = True
        return False

    pytesseract.pytesseract.tesseract_cmd = tesseract_cmd
    _configure_tesseract._configured = True
    return True


def _score_ocr_text(text: str) -> tuple[int, int, int]:
    compact_text = ''.join(ch for ch in text if ch.isalnum())
    return (len(compact_text), len(text.split()), len(text.strip()))


def _normalize_text(text: str) -> str:
    return ' '.join((text or '').split()).strip().lower()


def _prepare_image_for_ocr(image):
    from PIL import ImageFilter, ImageOps

    image = ImageOps.exif_transpose(image)

    if image.mode not in ['RGB', 'L']:
        image = image.convert('RGB')

    grayscale = ImageOps.grayscale(image)

    width, height = grayscale.size
    shortest_edge = max(1, min(width, height))
    if shortest_edge < 1400:
        scale = 1400 / shortest_edge
        resized = grayscale.resize((max(1, int(width * scale)), max(1, int(height * scale))))
    else:
        resized = grayscale

    enhanced = ImageOps.autocontrast(resized).filter(ImageFilter.SHARPEN)
    thresholded = enhanced.point(lambda value: 255 if value > 180 else 0)

    return [enhanced, thresholded]


def _extract_tesseract_text(image_path: str | None = None, image=None) -> str:
    if not _configure_tesseract():
        return ''

    import pytesseract
    from PIL import Image

    working_image = image
    if working_image is None and image_path:
        with Image.open(image_path) as opened_image:
            working_image = opened_image.copy()

    if working_image is None:
        return ''

    best_text = ''
    best_score = (0, 0, 0)
    configs = [
        '--oem 3 --psm 6 -c preserve_interword_spaces=1',
        '--oem 3 --psm 11 -c preserve_interword_spaces=1',
    ]

    for prepared_image in _prepare_image_for_ocr(working_image):
        for config in configs:
            try:
                text = pytesseract.image_to_string(prepared_image, config=config).strip()
            except Exception as exc:
                log.warning(f'Tesseract OCR failed: {exc}')
                continue

            score = _score_ocr_text(text)
            if score > best_score:
                best_text = text
                best_score = score

            if best_score[0] >= 80:
                break

    return best_text.strip()


class TesseractImageLoader:
    def __init__(self, file_path):
        self.file_path = file_path

    def load(self) -> list[Document]:
        text = _extract_tesseract_text(image_path=self.file_path)
        return [
            Document(
                page_content=text or '<No text content found>',
                metadata={'source': self.file_path, 'loader': 'tesseract_image'},
            )
        ]


class HybridPDFLoader:
    def __init__(self, file_path, extract_images=False, mode='page'):
        self.file_path = file_path
        self.extract_images = extract_images
        self.mode = mode

    def _open_pdf(self):
        try:
            import fitz
        except ImportError:
            log.warning('PyMuPDF is not installed. Falling back to PyPDFLoader for PDF text extraction.')
            return None

        return fitz.open(self.file_path)

    def _extract_native_page_text(self, page) -> str:
        blocks = page.get_text('blocks', sort=True)
        block_text = [
            str(block[4]).strip()
            for block in blocks
            if len(block) >= 5 and str(block[4]).strip()
        ]
        text = '\n\n'.join(block_text).strip()
        if text:
            return text
        return (page.get_text('text', sort=True) or '').strip()

    def _should_ocr_page(self, page, native_text: str) -> bool:
        if not self.extract_images:
            return False

        if not native_text.strip():
            return True

        image_count = len(page.get_images(full=True))
        if image_count == 0:
            return False

        text_word_count = len(native_text.split())
        text_char_count = len([ch for ch in native_text if not ch.isspace()])
        return text_word_count < 40 or text_char_count < 250

    def _render_page_for_ocr(self, page):
        import fitz
        from PIL import Image

        pixmap = page.get_pixmap(matrix=fitz.Matrix(2.5, 2.5), alpha=False)
        with Image.open(io.BytesIO(pixmap.tobytes('png'))) as rendered_image:
            return rendered_image.copy()

    def _build_page_documents(self) -> list[Document]:
        pdf = self._open_pdf()
        if pdf is None:
            return PyPDFLoader(
                self.file_path,
                extract_images=self.extract_images,
                mode=self.mode,
            ).load()

        docs = []
        try:
            total_pages = pdf.page_count
            for page in pdf:
                native_text = self._extract_native_page_text(page)
                page_text = native_text
                metadata = {
                    'source': self.file_path,
                    'page': page.number + 1,
                    'total_pages': total_pages,
                    'loader': 'hybrid_pdf',
                }

                if self._should_ocr_page(page, native_text):
                    try:
                        page_image = self._render_page_for_ocr(page)
                        ocr_text = _extract_tesseract_text(image=page_image)
                    except Exception as exc:
                        log.warning(f'Failed OCR fallback for PDF page {page.number + 1}: {exc}')
                        ocr_text = ''

                    normalized_page_text = _normalize_text(native_text)
                    normalized_ocr_text = _normalize_text(ocr_text)
                    if normalized_ocr_text and normalized_ocr_text not in normalized_page_text:
                        page_text = f'{native_text}\n\n{ocr_text}'.strip() if native_text else ocr_text
                        metadata['ocr'] = 'tesseract'

                docs.append(Document(page_content=page_text or '', metadata=metadata))
        finally:
            pdf.close()

        return docs

    def load(self) -> list[Document]:
        docs = self._build_page_documents()
        if self.mode != 'single':
            return docs

        combined_text = '\n\n'.join(
            [doc.page_content.strip() for doc in docs if (doc.page_content or '').strip()]
        ).strip()
        combined_metadata = {'source': self.file_path, 'loader': 'hybrid_pdf'}
        if any(doc.metadata.get('ocr') == 'tesseract' for doc in docs):
            combined_metadata['ocr'] = 'tesseract'
        return [Document(page_content=combined_text or '<No text content found>', metadata=combined_metadata)]


class Loader:
    def __init__(self, engine: str = '', **kwargs):
        self.engine = engine
        self.user = kwargs.get('user', None)
        self.kwargs = kwargs

    def load(self, filename: str, file_content_type: str, file_path: str) -> list[Document]:
        loader = self._get_loader(filename, file_content_type, file_path)
        docs = loader.load()
        return [Document(page_content=ftfy.fix_text(doc.page_content), metadata=doc.metadata) for doc in docs]

    async def aload(self, filename: str, file_content_type: str, file_path: str) -> list[Document]:
        """
        Async wrapper around `load`.

        Document loaders dispatched by `_get_loader` (PyMuPDF, Unstructured,
        python-docx, Tika, etc.) are uniformly synchronous and CPU/IO-bound.
        Calling `load` directly from an async handler would block the event
        loop for the entire parse — minutes for large PDFs. This offloads
        the work to a worker thread so the loop stays responsive.
        """
        return await asyncio.to_thread(self.load, filename, file_content_type, file_path)

    def _is_text_file(self, file_ext: str, file_content_type: str) -> bool:
        return file_ext in known_source_ext or (
            file_content_type
            and file_content_type.find('text/') >= 0
            # Avoid text/html files being detected as text
            and not file_content_type.find('html') >= 0
        )

    def _detect_text_encoding(self, file_path: str) -> str:
        """Detect the encoding of a text file with CJK-aware fallbacks.

        Langchain's ``TextLoader`` uses chardet internally when
        ``autodetect_encoding=True``, but chardet frequently misidentifies
        CJK encodings (e.g. GB18030 detected as GB2312 or even Cyrillic).
        This method replaces that by:

        1. Trying UTF-8 first (fast path for the vast majority of files).
        2. Using chardet as a *hint* to prioritise the right CJK codec
           family, but mapping subset names to their superset
           (e.g. GB2312 → gb18030).
        3. Validating that decoded text actually contains CJK characters,
           guarding against codecs that "succeed" but produce garbage.
        4. Falling back to latin-1 (always valid, ftfy fixes mojibake later).
        """
        try:
            with open(file_path, 'rb') as f:
                raw = f.read()
        except OSError:
            return 'utf-8'

        if not raw:
            return 'utf-8'

        # Fast path: most files are UTF-8
        try:
            raw.decode('utf-8')
            return 'utf-8'
        except UnicodeDecodeError:
            pass

        # Use chardet as a hint, not as ground truth
        import chardet

        detected = chardet.detect(raw)
        detected_enc = (detected.get('encoding') or '').lower().replace('-', '').replace('_', '')

        # Map chardet's detected encoding to the correct superset codec.
        # chardet often reports GB2312 for content that is actually GB18030;
        # GB18030 is a strict superset of both GB2312 and GBK.
        _ENC_FAMILY = {
            'gb2312': 'gb18030',
            'gb18030': 'gb18030',
            'gbk': 'gb18030',
            'big5': 'big5',
            'euckr': 'euc-kr',
            'eucjp': 'euc-jp',
            'iso2022jp': 'euc-jp',
            'shiftjis': 'shift_jis',
        }

        # Build priority list: chardet-hinted codec first, then remaining CJK
        base_order = ['gb18030', 'big5', 'euc-kr', 'euc-jp']
        hinted = _ENC_FAMILY.get(detected_enc)
        if hinted and hinted in base_order:
            ordered = [hinted] + [e for e in base_order if e != hinted]
        else:
            ordered = base_order

        for enc in ordered:
            try:
                text = raw.decode(enc)
                if text.strip() and self._has_cjk_characters(text):
                    log.info(
                        'Detected encoding %s for %s (chardet guessed %s)',
                        enc,
                        file_path,
                        detected.get('encoding'),
                    )
                    return enc
            except (UnicodeDecodeError, LookupError):
                continue

        # If chardet gave a non-CJK answer that isn't in our family map,
        # try it directly — it might be a valid Western encoding.
        chardet_encoding = detected.get('encoding')
        if chardet_encoding:
            try:
                raw.decode(chardet_encoding)
                log.info(
                    'Using chardet-detected encoding %s for %s',
                    chardet_encoding,
                    file_path,
                )
                return chardet_encoding
            except (UnicodeDecodeError, LookupError):
                pass

        # latin-1 is the ultimate fallback: every byte 0x00–0xFF is valid.
        # ftfy.fix_text() (applied downstream) repairs most mojibake that
        # results from treating Windows-1252 content as Latin-1.
        log.info('Falling back to latin-1 encoding for %s', file_path)
        return 'latin-1'

    @staticmethod
    def _has_cjk_characters(text: str, threshold: float = 0.05) -> bool:
        """Check if decoded text contains a meaningful proportion of CJK characters.

        This guards against codecs that technically "succeed" but decode the
        bytes into wrong Unicode codepoints (e.g. PUA chars, random symbols).
        A genuine CJK document should have at least ``threshold`` fraction of
        its non-whitespace characters in CJK Unicode blocks.
        """
        if not text:
            return False

        cjk_count = 0
        total = 0
        for ch in text:
            if ch.isspace():
                continue
            total += 1
            cp = ord(ch)
            if (
                0x4E00 <= cp <= 0x9FFF  # CJK Unified Ideographs
                or 0x3400 <= cp <= 0x4DBF  # CJK Extension A
                or 0x20000 <= cp <= 0x2A6DF  # CJK Extension B
                or 0x2A700 <= cp <= 0x2B73F  # CJK Extension C
                or 0x2B740 <= cp <= 0x2B81F  # CJK Extension D
                or 0xF900 <= cp <= 0xFAFF  # CJK Compatibility Ideographs
                or 0x3000 <= cp <= 0x303F  # CJK Symbols and Punctuation
                or 0x3040 <= cp <= 0x309F  # Hiragana
                or 0x30A0 <= cp <= 0x30FF  # Katakana
                or 0xAC00 <= cp <= 0xD7AF  # Hangul Syllables
                or 0xFF00 <= cp <= 0xFFEF  # Halfwidth and Fullwidth Forms
            ):
                cjk_count += 1

        if total == 0:
            return False

        return (cjk_count / total) >= threshold

    def _get_loader(self, filename: str, file_content_type: str, file_path: str):
        file_ext = filename.split('.')[-1].lower()

        if (
            self.engine == 'external'
            and self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_URL')
            and self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_API_KEY')
        ):
            loader = ExternalDocumentLoader(
                file_path=file_path,
                url=self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_URL'),
                api_key=self.kwargs.get('EXTERNAL_DOCUMENT_LOADER_API_KEY'),
                mime_type=file_content_type,
                user=self.user,
            )
        elif self.engine == 'tika' and self.kwargs.get('TIKA_SERVER_URL'):
            if self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            else:
                loader = TikaLoader(
                    url=self.kwargs.get('TIKA_SERVER_URL'),
                    file_path=file_path,
                    extract_images=self.kwargs.get('PDF_EXTRACT_IMAGES'),
                )
        elif (
            self.engine == 'datalab_marker'
            and self.kwargs.get('DATALAB_MARKER_API_KEY')
            and file_ext
            in [
                'pdf',
                'xls',
                'xlsx',
                'ods',
                'doc',
                'docx',
                'odt',
                'ppt',
                'pptx',
                'odp',
                'html',
                'epub',
                'png',
                'jpeg',
                'jpg',
                'webp',
                'gif',
                'tiff',
            ]
        ):
            api_base_url = self.kwargs.get('DATALAB_MARKER_API_BASE_URL', '')
            if not api_base_url or api_base_url.strip() == '':
                api_base_url = 'https://www.datalab.to/api/v1/marker'  # https://github.com/open-webui/open-webui/pull/16867#issuecomment-3218424349

            loader = DatalabMarkerLoader(
                file_path=file_path,
                api_key=self.kwargs['DATALAB_MARKER_API_KEY'],
                api_base_url=api_base_url,
                additional_config=self.kwargs.get('DATALAB_MARKER_ADDITIONAL_CONFIG'),
                use_llm=self.kwargs.get('DATALAB_MARKER_USE_LLM', False),
                skip_cache=self.kwargs.get('DATALAB_MARKER_SKIP_CACHE', False),
                force_ocr=self.kwargs.get('DATALAB_MARKER_FORCE_OCR', False),
                paginate=self.kwargs.get('DATALAB_MARKER_PAGINATE', False),
                strip_existing_ocr=self.kwargs.get('DATALAB_MARKER_STRIP_EXISTING_OCR', False),
                disable_image_extraction=self.kwargs.get('DATALAB_MARKER_DISABLE_IMAGE_EXTRACTION', False),
                format_lines=self.kwargs.get('DATALAB_MARKER_FORMAT_LINES', False),
                output_format=self.kwargs.get('DATALAB_MARKER_OUTPUT_FORMAT', 'markdown'),
            )
        elif self.engine == 'docling' and self.kwargs.get('DOCLING_SERVER_URL'):
            if self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            else:
                # Build params for DoclingLoader
                params = self.kwargs.get('DOCLING_PARAMS', {})
                if not isinstance(params, dict):
                    try:
                        params = json.loads(params)
                    except json.JSONDecodeError:
                        log.error('Invalid DOCLING_PARAMS format, expected JSON object')
                        params = {}

                loader = DoclingLoader(
                    url=self.kwargs.get('DOCLING_SERVER_URL'),
                    api_key=self.kwargs.get('DOCLING_API_KEY', None),
                    file_path=file_path,
                    mime_type=file_content_type,
                    params=params,
                )
        elif (
            self.engine == 'document_intelligence'
            and self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT') != ''
            and (
                file_ext in ['pdf', 'docx', 'ppt', 'pptx']
                or file_content_type
                in [
                    'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
                    'application/vnd.ms-powerpoint',
                    'application/vnd.openxmlformats-officedocument.presentationml.presentation',
                ]
            )
        ):
            if self.kwargs.get('DOCUMENT_INTELLIGENCE_KEY') != '':
                loader = AzureAIDocumentIntelligenceLoader(
                    file_path=file_path,
                    api_endpoint=self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT'),
                    api_key=self.kwargs.get('DOCUMENT_INTELLIGENCE_KEY'),
                    api_model=self.kwargs.get('DOCUMENT_INTELLIGENCE_MODEL'),
                )
            else:
                loader = AzureAIDocumentIntelligenceLoader(
                    file_path=file_path,
                    api_endpoint=self.kwargs.get('DOCUMENT_INTELLIGENCE_ENDPOINT'),
                    azure_credential=DefaultAzureCredential(),
                    api_model=self.kwargs.get('DOCUMENT_INTELLIGENCE_MODEL'),
                )
        elif self.engine == 'mineru' and file_ext in self.kwargs.get('MINERU_FILE_EXTENSIONS', ['pdf']):
            mineru_timeout = self.kwargs.get('MINERU_API_TIMEOUT', 300)
            if mineru_timeout:
                try:
                    mineru_timeout = int(mineru_timeout)
                except ValueError:
                    mineru_timeout = 300

            loader = MinerULoader(
                file_path=file_path,
                api_mode=self.kwargs.get('MINERU_API_MODE', 'local'),
                api_url=self.kwargs.get('MINERU_API_URL', 'http://localhost:8000'),
                api_key=self.kwargs.get('MINERU_API_KEY', ''),
                params=self.kwargs.get('MINERU_PARAMS', {}),
                timeout=mineru_timeout,
            )
        elif (
            self.engine == 'mistral_ocr'
            and self.kwargs.get('MISTRAL_OCR_API_KEY') != ''
            and file_ext in ['pdf']  # Mistral OCR currently only supports PDF and images
        ):
            loader = MistralLoader(
                base_url=self.kwargs.get('MISTRAL_OCR_API_BASE_URL'),
                api_key=self.kwargs.get('MISTRAL_OCR_API_KEY'),
                file_path=file_path,
            )
        elif self.engine == 'paddleocr_vl' and self.kwargs.get('PADDLEOCR_VL_TOKEN') != '':
            loader = PaddleOCRVLLoader(
                api_url=self.kwargs.get('PADDLEOCR_VL_BASE_URL'),
                token=self.kwargs.get('PADDLEOCR_VL_TOKEN'),
                file_path=file_path,
            )
        else:
            if file_ext == 'pdf':
                loader = HybridPDFLoader(
                    file_path,
                    extract_images=self.kwargs.get('PDF_EXTRACT_IMAGES'),
                    mode=self.kwargs.get('PDF_LOADER_MODE', 'page'),
                )
            elif file_ext == 'csv':
                loader = CSVLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_ext == 'rst':
                try:
                    from langchain_community.document_loaders import UnstructuredRSTLoader

                    loader = UnstructuredRSTLoader(file_path, mode='elements')
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to plain text loading for .rst file. '
                        'Install it with: pip install unstructured'
                    )
                    loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_ext == 'xml':
                try:
                    from langchain_community.document_loaders import UnstructuredXMLLoader

                    loader = UnstructuredXMLLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to plain text loading for .xml file. '
                        'Install it with: pip install unstructured'
                    )
                    loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_ext in ['htm', 'html']:
                loader = BSHTMLLoader(file_path, open_encoding='unicode_escape')
            elif file_ext == 'md':
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            elif file_content_type == 'application/epub+zip':
                try:
                    from langchain_community.document_loaders import UnstructuredEPubLoader

                    loader = UnstructuredEPubLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .epub files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif (
                file_content_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'
                or file_ext == 'docx'
            ):
                loader = Docx2txtLoader(file_path)
            elif file_ext == 'doc' or file_content_type == 'application/msword':
                try:
                    from langchain_community.document_loaders import UnstructuredWordDocumentLoader

                    loader = UnstructuredWordDocumentLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .doc files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif file_content_type in [
                'application/vnd.ms-excel',
                'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            ] or file_ext in ['xls', 'xlsx']:
                try:
                    from langchain_community.document_loaders import UnstructuredExcelLoader

                    loader = UnstructuredExcelLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to pandas for Excel file loading. '
                        'Install unstructured for better results: pip install unstructured'
                    )
                    loader = ExcelLoader(file_path)
            elif file_content_type in [
                'application/vnd.ms-powerpoint',
                'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            ] or file_ext in ['ppt', 'pptx']:
                try:
                    from langchain_community.document_loaders import UnstructuredPowerPointLoader

                    loader = UnstructuredPowerPointLoader(file_path)
                except ImportError:
                    log.warning(
                        "The 'unstructured' package is not installed. "
                        'Falling back to python-pptx for PowerPoint file loading. '
                        'Install unstructured for better results: pip install unstructured'
                    )
                    loader = PptxLoader(file_path)
            elif file_ext == 'msg':
                loader = OutlookMessageLoader(file_path)
            elif file_ext == 'odt':
                try:
                    from langchain_community.document_loaders import UnstructuredODTLoader

                    loader = UnstructuredODTLoader(file_path)
                except ImportError:
                    raise ValueError(
                        "Processing .odt files requires the 'unstructured' package. "
                        'Install it with: pip install unstructured'
                    )
            elif file_ext in image_file_ext or (file_content_type and file_content_type.startswith('image/')):
                loader = TesseractImageLoader(file_path)
            elif self._is_text_file(file_ext, file_content_type):
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))
            else:
                loader = TextLoader(file_path, encoding=self._detect_text_encoding(file_path))

        return loader
